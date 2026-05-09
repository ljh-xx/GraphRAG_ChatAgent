from __future__ import annotations

import hashlib
import io
import re
from pathlib import Path
from typing import Callable

import fitz
from docx import Document as DocxDocument
from PIL import Image
from pptx import Presentation

from .models import ParsedDocument

try:
    from rapidocr_onnxruntime import RapidOCR
except Exception:  # pragma: no cover - optional dependency runtime guard
    RapidOCR = None


SUPPORTED_SUFFIXES = {".pdf", ".docx", ".doc", ".pptx", ".png", ".jpg", ".jpeg", ".webp", ".md", ".markdown", ".txt", ".xlsx", ".xls"}


class DocumentParser:
    def __init__(self, progress_callback: Callable[[str], None] | None = None) -> None:
        self.ocr_engine = RapidOCR() if RapidOCR else None
        self.progress_callback = progress_callback

    def parse_path(self, path: str | Path) -> list[ParsedDocument]:
        file_path = Path(path)
        suffix = file_path.suffix.lower()
        self._log(f"[parse] 开始处理文件: {file_path.name} ({suffix})")

        if suffix == ".pdf":
            return self._parse_pdf(file_path)
        if suffix == ".docx":
            return self._parse_docx(file_path)
        if suffix == ".doc":
            return self._parse_doc_with_word(file_path)
        if suffix == ".pptx":
            return self._parse_pptx(file_path)
        if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
            return self._parse_image(file_path)
        if suffix in {".md", ".markdown"}:
            return self._parse_text(file_path, "markdown")
        if suffix == ".txt":
            return self._parse_text(file_path, "text")
        if suffix in {".xlsx", ".xls"}:
            return self._parse_excel(file_path)
        raise ValueError(f"Unsupported file type: {file_path.suffix}")

    def parse_directory(self, directory: str | Path) -> list[ParsedDocument]:
        root = Path(directory)
        results: list[ParsedDocument] = []
        files = [path for path in sorted(root.rglob("*")) if path.is_file() and path.suffix.lower() in SUPPORTED_SUFFIXES]
        self._log(f"[parse] 在目录中发现 {len(files)} 个可处理文件: {root}")
        for index, path in enumerate(files, start=1):
            self._log(f"[parse] 文件进度 {index}/{len(files)}: {path.name}")
            if path.is_file() and path.suffix.lower() in SUPPORTED_SUFFIXES:
                results.extend(self.parse_path(path))
        self._log(f"[parse] 文档解析完成，共生成 {len(results)} 条文档记录")
        return results

    def _parse_pdf(self, path: Path) -> list[ParsedDocument]:
        doc = fitz.open(path)
        doc_id = self._file_id(path)
        self._log(f"[parse][pdf] {path.name} 共 {len(doc)} 页")

        full_text_with_pages = []
        for page_index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if len(text) < 40:
                self._log(f"[parse][pdf] 第 {page_index} 页文本较少，尝试 OCR")
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                image_bytes = pix.tobytes("png")
                text = self._ocr_image_bytes(image_bytes)
            
            if text.strip():
                full_text_with_pages.append((page_index, text))

        if not full_text_with_pages:
            self._log(f"[parse][pdf] {path.name} 无有效文本")
            return []

        # 尝试通过正则表达式识别章节
        # 常见论文标题模式: 1. Introduction, II. METHODS, Abstract, References
        section_pattern = re.compile(r"^(\d+\.[\s\t]+|[I|V|X]+\.[\s\t]+|Abstract|References|Conclusion|Introduction|Methodology|Related Work)", re.IGNORECASE | re.MULTILINE)
        
        combined_text = "\n".join([t for _, t in full_text_with_pages])
        
        # 简单的章节识别：如果匹配到上述模式，且该行字符数不多，认为是标题
        records: list[ParsedDocument] = []
        lines = combined_text.splitlines()
        current_section_text = []
        current_section_title = "Abstract/Intro"
        current_page = full_text_with_pages[0][0]

        for line in lines:
            stripped = line.strip()
            if section_pattern.match(stripped) and len(stripped) < 100:
                if current_section_text:
                    records.append(ParsedDocument(
                        source_path=str(path),
                        doc_id=doc_id,
                        title=f"{path.stem} - {current_section_title}",
                        text="\n".join(current_section_text).strip(),
                        modality="pdf",
                        page_number=current_page,
                        extra_meta={"section": current_section_title}
                    ))
                current_section_title = stripped
                current_section_text = [line]
            else:
                current_section_text.append(line)

        # 最后一节
        if current_section_text:
            records.append(ParsedDocument(
                source_path=str(path),
                doc_id=doc_id,
                title=f"{path.stem} - {current_section_title}",
                text="\n".join(current_section_text).strip(),
                modality="pdf",
                page_number=current_page,
                extra_meta={"section": current_section_title}
            ))

        self._log(f"[parse][pdf] {path.name} 解析完成，识别到 {len(records)} 个章节")
        return records

    def _parse_docx(self, path: Path) -> list[ParsedDocument]:
        doc = DocxDocument(path)
        text = "\n".join([para.text.strip() for para in doc.paragraphs if para.text.strip()])
        self._log(f"[parse][docx] {path.name} 提取到 {len(text)} 个字符")
        return self._single_record(path, text, "docx")

    def _parse_doc_with_word(self, path: Path) -> list[ParsedDocument]:
        try:
            import win32com.client  # type: ignore
        except Exception as exc:  # pragma: no cover - only for Windows desktop
            raise RuntimeError(f"Reading .doc requires pywin32 and Word. File: {path}") from exc

        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        document = word.Documents.Open(str(path))
        try:
            text = document.Content.Text
        finally:
            document.Close(False)
            word.Quit()
        self._log(f"[parse][doc] {path.name} 提取到 {len(text)} 个字符")
        return self._single_record(path, text, "doc")

    def _parse_pptx(self, path: Path) -> list[ParsedDocument]:
        prs = Presentation(path)
        records: list[ParsedDocument] = []
        doc_id = self._file_id(path)
        self._log(f"[parse][pptx] {path.name} 共 {len(prs.slides)} 页")
        for slide_idx, slide in enumerate(prs.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text.strip())
            text = "\n".join([line for line in lines if line])
            if text:
                records.append(
                    ParsedDocument(
                        source_path=str(path),
                        doc_id=doc_id,
                        title=path.stem,
                        text=text,
                        modality="pptx_slide",
                        page_number=slide_idx,
                    )
                )
        self._log(f"[parse][pptx] {path.name} 解析完成，保留 {len(records)} 页")
        return records

    def _parse_image(self, path: Path) -> list[ParsedDocument]:
        with path.open("rb") as f:
            image_bytes = f.read()
        self._log(f"[parse][image] {path.name} 开始 OCR")
        text = self._ocr_image_bytes(image_bytes)
        self._log(f"[parse][image] {path.name} OCR 提取到 {len(text)} 个字符")
        return self._single_record(path, text, "image")

    def _parse_text(self, path: Path, modality: str = "text") -> list[ParsedDocument]:
        try:
            with path.open("r", encoding="utf-8") as f:
                text = f.read()
        except UnicodeDecodeError:
            # 如果 utf-8 失败，尝试 gbk (针对某些 Windows 环境下的文本文件)
            with path.open("r", encoding="gbk") as f:
                text = f.read()
        
        self._log(f"[parse][{modality}] {path.name} 提取到 {len(text)} 个字符")
        return self._single_record(path, text, modality)

    def _parse_excel(self, path: Path) -> list[ParsedDocument]:
        import pandas as pd
        # 读取所有工作表
        all_sheets = pd.read_excel(path, sheet_name=None)
        records: list[ParsedDocument] = []
        doc_id = self._file_id(path)
        
        for sheet_name, df in all_sheets.items():
            self._log(f"[parse][excel] 正在处理工作表: {sheet_name}")
            
            # 尝试匹配可能的列名（不区分大小写）
            q_col = next((c for c in df.columns if str(c).lower() in ["question", "问题", "q"]), None)
            a_col = next((c for c in df.columns if str(c).lower() in ["answer", "答案", "a"]), None)
            category_col = next((c for c in df.columns if str(c).lower() in ["category", "分类", "类型", "主题", "topic"]), None)

            if q_col is not None and a_col is not None:
                self._log(f"[parse][excel] {path.name}[{sheet_name}] 识别到问答列: {q_col}, {a_col}")
                
                qa_buffer = []
                current_buffer_len = 0
                max_buffer_len = 1000 # 针对短 Q&A 进行聚合，减少后续 LLM 抽取压力
                
                for idx, row in df.iterrows():
                    q = str(row[q_col]).strip() if pd.notna(row[q_col]) else ""
                    a = str(row[a_col]).strip() if pd.notna(row[a_col]) else ""
                    cat = str(row[category_col]).strip() if category_col and pd.notna(row[category_col]) else sheet_name
                    
                    if q or a:
                        qa_text = f"问题: {q}\n答案: {a}"
                        qa_buffer.append(qa_text)
                        current_buffer_len += len(qa_text)
                        
                        if current_buffer_len >= max_buffer_len:
                            combined_text = "\n---\n".join(qa_buffer)
                            records.append(ParsedDocument(
                                source_path=str(path),
                                doc_id=doc_id,
                                title=f"{path.stem} - {sheet_name}",
                                text=combined_text,
                                modality="excel_qa",
                                page_number=idx + 1,
                                extra_meta={"sheet": sheet_name, "category": cat, "is_grouped": True}
                            ))
                            qa_buffer = []
                            current_buffer_len = 0
                
                # 处理剩余 buffer
                if qa_buffer:
                    combined_text = "\n---\n".join(qa_buffer)
                    records.append(ParsedDocument(
                        source_path=str(path),
                        doc_id=doc_id,
                        title=f"{path.stem} - {sheet_name}",
                        text=combined_text,
                        modality="excel_qa",
                        page_number=len(df),
                        extra_meta={"sheet": sheet_name, "is_grouped": True}
                    ))
            else:
                self._log(f"[parse][excel] {path.name}[{sheet_name}] 未能找到明确的问答列，按行/分类合并")
                if category_col:
                    # 如果有分类列，按分类合并
                    for cat_name, group in df.groupby(category_col):
                        text_parts = []
                        for _, row in group.iterrows():
                            row_text = " | ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val) and col != category_col])
                            text_parts.append(row_text)
                        
                        combined_text = f"分类: {cat_name}\n" + "\n".join(text_parts)
                        records.append(ParsedDocument(
                            source_path=str(path),
                            doc_id=doc_id,
                            title=f"{path.stem} - {sheet_name} - {cat_name}",
                            text=combined_text,
                            modality="excel_row",
                            page_number=0, # 分组数据无固定行号
                            extra_meta={"sheet": sheet_name, "category": cat_name}
                        ))
                else:
                    # 无分类列，按行处理
                    for idx, row in df.iterrows():
                        text = " ".join([str(v) for v in row.values if pd.notna(v)])
                        if text.strip():
                            records.append(ParsedDocument(
                                source_path=str(path),
                                doc_id=doc_id,
                                title=f"{path.stem} - {sheet_name}",
                                text=text,
                                modality="excel_row",
                                page_number=idx + 1,
                                extra_meta={"sheet": sheet_name}
                            ))
        
        self._log(f"[parse][excel] {path.name} 解析完成，共提取 {len(records)} 条记录")
        return records

    def _single_record(self, path: Path, text: str, modality: str) -> list[ParsedDocument]:
        cleaned = text.strip()
        if not cleaned:
            return []
        return [
            ParsedDocument(
                source_path=str(path),
                doc_id=self._file_id(path),
                title=path.stem,
                text=cleaned,
                modality=modality,
            )
        ]

    def _ocr_image_bytes(self, image_bytes: bytes) -> str:
        if not self.ocr_engine:
            self._log("[parse][ocr] 未检测到 OCR 引擎，返回空文本")
            return ""
        image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        result, _ = self.ocr_engine(image)
        if not result:
            self._log("[parse][ocr] OCR 未识别到文本")
            return ""
        return "\n".join([item[1] for item in result if len(item) > 1 and item[1]])

    @staticmethod
    def _file_id(path: Path) -> str:
        digest = hashlib.md5(str(path.resolve()).encode("utf-8")).hexdigest()
        return f"doc_{digest}"

    def _log(self, message: str) -> None:
        if self.progress_callback:
            self.progress_callback(message)
